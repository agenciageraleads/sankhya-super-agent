import os
from docx import Document
from pptx import Presentation
import pandas as pd
import json

def extract_docx(filepath):
    doc = Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_pptx(filepath):
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_xlsx(filepath):
    df = pd.read_excel(filepath)
    return df.to_markdown()

training_root = "/Users/Lucas-Lenovo/Sankhya-Super-Agente/sankhya-agent/mcp_server/domains/procurement/training"
results = {}

# Process Docs
docs_dir = os.path.join(training_root, "docs")
for f in os.listdir(docs_dir):
    path = os.path.join(docs_dir, f)
    if f.endswith(".docx"):
        results[f] = extract_docx(path)
    elif f.endswith(".xlsx"):
        results[f] = extract_xlsx(path)

# Process Slides
slides_dir = os.path.join(training_root, "slides")
for f in os.listdir(slides_dir):
    path = os.path.join(slides_dir, f)
    if f.endswith(".pptx"):
        results[f] = extract_pptx(path)

# Save results for agent to read
with open(os.path.join(training_root, "extracted_content.json"), "w", encoding="utf-8") as f:
    json.dump(results, f, indent=2, ensure_ascii=False)

print("Extraction complete. JSON saved.")
